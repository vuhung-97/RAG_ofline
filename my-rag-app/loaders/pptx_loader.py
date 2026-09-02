from pptx import Presentation
from loaders.base import BaseLoader

class PptxLoader(BaseLoader):
    """SRP: Trích xuất văn bản từ PowerPoint (.pptx)."""

    def load(self, file_path: str) -> str:
        prs = Presentation(file_path)
        text_parts = []

        for slide_idx, slide in enumerate(prs.slides, 1):
            text_parts.append(f"--- Slide {slide_idx} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())

        return "\n".join(text_parts)
